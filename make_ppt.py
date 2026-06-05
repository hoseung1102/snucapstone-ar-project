#!/usr/bin/env python3
"""Eagle Eye 발표자료 생성 — 기존 2차 미팅 PPT + 비교분석 + 하드웨어 + 메모리"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── 색상 ─────────────────────────────────────────────────────────────
BG     = RGBColor(0x0B, 0x0D, 0x10)
CARD   = RGBColor(0x11, 0x14, 0x1A)
CARD_W = RGBColor(0x1C, 0x1A, 0x18)
ORANGE = RGBColor(0xE3, 0xA2, 0x4A)
WHITE  = RGBColor(0xE8, 0xE6, 0xE1)
GRAY   = RGBColor(0x8A, 0x8B, 0x87)
GREEN  = RGBColor(0x4A, 0xD6, 0x6D)
RED    = RGBColor(0xFF, 0x6B, 0x6B)
HILITE = RGBColor(0x14, 0x1E, 0x30)

FONT = "Arial"
SW   = Emu(18288000)
SH   = Emu(10287000)
LM   = Emu(952500)
RM   = Emu(952500)
CW   = SW - LM - RM   # 16383000

# ── 헬퍼 ─────────────────────────────────────────────────────────────

def set_cell_fill(cell, color):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for el in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(el)
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    sc = etree.SubElement(sf, qn('a:srgbClr'))
    sc.set('val', f'{color[0]:02X}{color[1]:02X}{color[2]:02X}')

def cell_text(cell, text, size=13, bold=False, color=WHITE, align=PP_ALIGN.CENTER):
    tf = cell.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

def box(slide, text, x, y, w, h, size=16, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in text.split('\n'):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color

def rect(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = fill
    return s

def bg(slide):
    rect(slide, Emu(0), Emu(0), SW, SH, BG)

def label(slide, num, text):
    box(slide, f"{num:02d} · {text}",
        LM, Emu(381000), Emu(5000000), Emu(250000),
        size=11, color=ORANGE)

def title(slide, text):
    box(slide, text, LM, Emu(1050000), CW, Emu(900000), size=75, color=WHITE)

def sub(slide, text, size=22):
    box(slide, text, LM, Emu(2100000), CW, Emu(550000), size=size, color=GRAY)

def move_slide_to_end(prs, idx):
    lst = prs.slides._sldIdLst
    el = list(lst)[idx]
    lst.remove(el)
    lst.append(el)

# ── 기존 PPT 로드 ────────────────────────────────────────────────────
SRC = "/Users/choehoseung/Desktop/Hoseung/SNU/2026-1/학제간 캡스톤 설계/2차 프로그레스 미팅_3팀.pptx"
DST = "/Users/choehoseung/Desktop/Hoseung/SNU/2026-1/학제간 캡스톤 설계/Eagle_Eye_발표자료.pptx"

prs = Presentation(SRC)
layout = prs.slide_layouts[0]

# ════════════════════════════════════════════════════════════════════
# 슬라이드 3: RELATED WORK — 기존 연구 비교
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(layout)
bg(sl)
label(sl, 3, "RELATED WORK")
title(sl, "PRIOR ART")
sub(sl, "셋 다 만족하는 건 Eagle Eye뿐.")

HDR  = ["", "온디바이스", "자동 트리거", "마커 없음"]
ROWS = [
    ("TinyissimoYOLO", "✅",           "❌  always-on",  "✅"),
    ("ARShopping",     "❌",           "❌",             "❌"),
    ("MobileCLIP",     "✅",           "—  모델만",      "—"),
    ("SUPERGLASSES",   "❌  클라우드VLM","❌  음성 호출", "✅"),
    ("Meta Ray-Ban",   "❌  클라우드", "❌  Hey Meta",   "✅"),
    ("Eagle Eye  ★",  "✅",           "✅",             "✅"),
]

tbl = sl.shapes.add_table(
    len(ROWS) + 1, 4,
    LM, Emu(2800000), CW, Emu(6700000)
).table
tbl.columns[0].width = Emu(5100000)
tbl.columns[1].width = Emu(3761000)
tbl.columns[2].width = Emu(3900000)
tbl.columns[3].width = Emu(3622000)

HDR_BG = RGBColor(0x18, 0x0E, 0x02)
for ci, h in enumerate(HDR):
    c = tbl.cell(0, ci)
    set_cell_fill(c, HDR_BG)
    cell_text(c, h, size=14, bold=True, color=ORANGE if ci > 0 else GRAY)

for ri, row in enumerate(ROWS):
    eagle = ri == len(ROWS) - 1
    rbg = HILITE if eagle else CARD
    for ci, val in enumerate(row):
        c = tbl.cell(ri + 1, ci)
        set_cell_fill(c, rbg)
        clr = (GREEN  if val.startswith("✅") else
               RED    if val.startswith("❌") else
               GRAY   if val == "—" else
               ORANGE if eagle else WHITE)
        al = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        cell_text(c, val, size=14 if ci == 0 else 13, bold=eagle, color=clr, align=al)

box(sl, "→ 세 조건(온디바이스 · 자동 트리거 · 마커 없음)을 동시에 만족하는 시스템은 Eagle Eye가 유일합니다.",
    LM, Emu(9700000), CW, Emu(400000), size=13, color=GRAY)

# ════════════════════════════════════════════════════════════════════
# 슬라이드 4: TARGET HARDWARE — RayNeo X3 Pro 선정 이유
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(layout)
bg(sl)
label(sl, 4, "HARDWARE")
title(sl, "TARGET HARDWARE")
sub(sl, "AR1 NPU + SLAM + 글라스 폼팩터 — 현재 구매 가능한 유일한 옵션.")

# 요구조건 카드 5개
REQ = [
    ("Snapdragon\nAR1 Gen 1",    "NPU 탑재\nYOLO+CLIP 온디바이스"),
    ("내장 SLAM\n(6DoF)",         "Spatial Anchor\n필수 조건"),
    ("양안 디스플레이\nMicroLED 6,000nit", "실외 매장\n환경 대응"),
    ("글라스 폼팩터\n76g",         "일상 착용 가능\nHMD 아님"),
    ("Android SDK\n(AOSP)",       "개발 진입장벽\n낮음"),
]
CRD_W = Emu(3076600)
CRD_H = Emu(2500000)
CRD_GAP = Emu(220750)
CRD_Y = Emu(2900000)

for i, (t, d) in enumerate(REQ):
    cx = LM + i * (CRD_W + CRD_GAP)
    rect(sl, cx, CRD_Y, CRD_W, CRD_H, CARD)
    box(sl, "✅", cx + Emu(180000), CRD_Y + Emu(160000),
        Emu(400000), Emu(450000), size=20, color=GREEN)
    box(sl, t, cx + Emu(180000), CRD_Y + Emu(580000),
        CRD_W - Emu(360000), Emu(800000), size=13, color=WHITE, bold=True)
    box(sl, d, cx + Emu(180000), CRD_Y + Emu(1480000),
        CRD_W - Emu(360000), Emu(900000), size=12, color=GRAY)

# 하단 2열: 대안 대비 / 한계 인정
LY = CRD_Y + CRD_H + Emu(280000)
LH = Emu(4100000)
LW = Emu(7700000)
RX = LM + LW + Emu(300000)
RW = CW - LW - Emu(300000)

rect(sl, LM, LY, LW, LH, CARD_W)
box(sl, "대안 대비 우위", LM + Emu(250000), LY + Emu(180000),
    LW - Emu(500000), Emu(380000), size=15, color=ORANGE, bold=True)

ALTS = [
    ("XREAL Air 2 Ultra",   "Spatial Anchor 검증됐지만 단종 임박"),
    ("XREAL Project Aura",  "2026 출시 예정, 현재 구매 불가"),
    ("Meta Ray-Ban",        "디스플레이 없음 — 오디오+카메라만, AR 불가"),
    ("HoloLens / Vision Pro","HMD 폼팩터, 일상 착용 부적합"),
]
for j, (nm, rs) in enumerate(ALTS):
    py = LY + Emu(730000) + j * Emu(780000)
    box(sl, nm, LM + Emu(250000), py,
        LW - Emu(500000), Emu(280000), size=13, color=RED, bold=True)
    box(sl, f"→ {rs}", LM + Emu(250000), py + Emu(280000),
        LW - Emu(500000), Emu(330000), size=12, color=GRAY)

rect(sl, RX, LY, RW, LH, CARD)
box(sl, "한계 인정", RX + Emu(250000), LY + Emu(180000),
    RW - Emu(500000), Emu(380000), size=15, color=ORANGE, bold=True)

LIMS = [
    ("Spatial Anchor API 공개 여부",
     "RayNeo 개발자 답변 대기 중. SDK 문서에 미공개 API 존재 가능성 있음."),
    ("Eye-tracking 없음",
     "head-gaze로 대체 → 트리거 정밀도 trade-off. 실사용 검증 필요."),
]
for j, (nm, ds) in enumerate(LIMS):
    py = LY + Emu(730000) + j * Emu(1300000)
    box(sl, nm, RX + Emu(250000), py,
        RW - Emu(500000), Emu(300000), size=14, color=ORANGE, bold=True)
    box(sl, ds, RX + Emu(250000), py + Emu(330000),
        RW - Emu(500000), Emu(700000), size=12, color=GRAY)

# ════════════════════════════════════════════════════════════════════
# 슬라이드 5: MEMORY FOOTPRINT — RAM 분석
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(layout)
bg(sl)
label(sl, 5, "MEMORY")
title(sl, "MEMORY FOOTPRINT")
sub(sl, "RayNeo X3 Pro 4GB RAM — Eagle Eye 앱 여유 ~1GB 이상.")

MEM = [
    ("YOLOv11n (INT8)",                  "~10 MB",      "온디바이스 추론"),
    ("MobileCLIP-S0 (INT8)",             "~50 MB",      "상품 식별 모델"),
    ("Vector DB\n(1만 상품, 512-dim)",   "~20 MB",      "float32 기준"),
    ("SLAM 엔진 (시스템)",                "~500MB – 1GB","OS 레벨 공유"),
    ("OS + 기타",                        "~1.5 GB",     "Android 기준"),
    ("Eagle Eye 앱 여유",                "~1 GB 이상",  "모델+DB 합계 <100MB"),
]
TW = Emu(11200000)
tbl2 = sl.shapes.add_table(
    len(MEM) + 1, 3,
    LM, Emu(2800000), TW, Emu(6500000)
).table
tbl2.columns[0].width = Emu(5000000)
tbl2.columns[1].width = Emu(3000000)
tbl2.columns[2].width = Emu(3200000)

for ci, h in enumerate(["컴포넌트", "예상 사용량", "비고"]):
    c = tbl2.cell(0, ci)
    set_cell_fill(c, RGBColor(0x18, 0x0E, 0x02))
    cell_text(c, h, size=14, bold=True, color=ORANGE)

for ri, (comp, sz, note) in enumerate(MEM):
    last = ri == len(MEM) - 1
    rbg = HILITE if last else CARD
    for ci, (val, clr, al) in enumerate([
        (comp, ORANGE if last else WHITE, PP_ALIGN.LEFT),
        (sz,   GREEN  if last else WHITE, PP_ALIGN.CENTER),
        (note, GRAY,                      PP_ALIGN.LEFT),
    ]):
        c = tbl2.cell(ri + 1, ci)
        set_cell_fill(c, rbg)
        cell_text(c, val, size=13, bold=last, color=clr, align=al)

# 결론 박스 (오른쪽)
BX = LM + TW + Emu(400000)
BW = CW - TW - Emu(400000)
BH = Emu(6500000)
rect(sl, BX, Emu(2800000), BW, BH, CARD_W)

CONC = [
    (Emu(2950000), "결론",                    18, ORANGE, True),
    (Emu(3450000), "모델+DB 합계",             14, GRAY,   False),
    (Emu(3800000), "< 100 MB",                22, GREEN,  True),
    (Emu(4350000), "4GB RAM으로 충분.",        14, WHITE,  False),
    (Emu(5000000), "확장 시나리오",            13, ORANGE, True),
    (Emu(5350000), "10만 상품 → ~200MB",       12, GRAY,   False),
    (Emu(5650000), "MobileCLIP-S2 → ~150MB",   12, GRAY,   False),
    (Emu(6200000), "→ 모두 여유 있음",          13, GREEN,  True),
]
for (cy, text, sz, clr, bld) in CONC:
    box(sl, text, BX + Emu(280000), cy, BW - Emu(560000), Emu(450000),
        size=sz, color=clr, bold=bld)

# ════════════════════════════════════════════════════════════════════
# NEXT 슬라이드를 끝으로 이동 (03→06 으로)
# 현재 순서: [title(0), pipeline(1), next(2), related(3), hw(4), mem(5)]
# 목표:      [title(0), pipeline(1), related(2), hw(3), mem(4), next(5)]
# ════════════════════════════════════════════════════════════════════
move_slide_to_end(prs, 2)

# 섹션 레이블 "03 · NEXT" → "06 · NEXT" 업데이트
for shape in prs.slides[5].shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '03 · NEXT' in run.text:
                    run.text = run.text.replace('03 · NEXT', '06 · NEXT')

prs.save(DST)
print(f"완료: {DST}")
print(f"총 슬라이드: {len(prs.slides)}장")
